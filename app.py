import streamlit as st
import openai
from langchain.chains.question_answering import load_qa_chain
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAI
import pypdf
import os
from io import BytesIO
import docx  # สำหรับไฟล์ .docx
from pptx import Presentation  # สำหรับไฟล์ .pptx
import pandas as pd  # สำหรับไฟล์ .xlsx
from langchain.memory import ConversationBufferMemory
from pymongo import MongoClient
import uuid
from datetime import datetime, UTC

# --- ฟังก์ชันหลักในการทำงาน ---
def get_text_from_file(file):
    """
    ฟังก์ชันสำหรับอ่านข้อความจากไฟล์ที่อัปโหลด (PDF, TXT, DOCX, PPTX, XLSX, CSV)
    """
    text = ""
    file_extension = os.path.splitext(file.name)[1]
    if file_extension == ".pdf":
        try:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาดในการอ่านไฟล์ PDF: {e}")
            return None
    elif file_extension == ".txt":
        text = file.getvalue().decode("utf-8")
    elif file_extension == ".docx":
        try:
            document = docx.Document(file)
            for para in document.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาดในการอ่านไฟล์ DOCX: {e}")
            return None
    elif file_extension == ".pptx":
        try:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาดในการอ่านไฟล์ PPTX: {e}")
            return None
    elif file_extension == ".xlsx":
        try:
            xls = pd.ExcelFile(file)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name, header=0)
                # แปลง DataFrame ทั้งหมดเป็นข้อความ
                text += f"--- Sheet: {sheet_name} ---\n"
                text += df.to_string(index=False) + "\n\n"
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาดในการอ่านไฟล์ XLSX: {e}")
            return None
    elif file_extension == ".csv":
        try:
            df = pd.read_csv(file)
            text += df.to_string(index=False) + "\n"
        except Exception as e:
            st.error(f"เกิดข้อผิดพลาดในการอ่านไฟล์ CSV: {e}")
            return None
    else:
        st.warning("ไม่รองรับชนิดของไฟล์นี้ กรุณาอัปโหลดไฟล์ PDF, TXT, DOCX, PPTX, XLSX หรือ CSV", icon="⚠️")
        return None
    return text

def get_text_chunks(raw_text):
    """
    ฟังก์ชันสำหรับแบ่งข้อความยาวๆ ออกเป็นส่วนย่อย (chunks)
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(raw_text)
    return chunks

def get_vector_store(text_chunks, api_key):
    """
    ฟังก์ชันสำหรับแปลง chunks ของข้อความเป็น vector embeddings และเก็บใน FAISS vector store
    """
    try:
        embeddings = OpenAIEmbeddings(openai_api_key=api_key)
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        return vector_store
    except Exception as e:
        st.error(f"เกิดข้อผิดพลาดระหว่างสร้าง Vector Store: {e}")
        st.info("อาจเกิดจาก API Key ไม่ถูกต้อง หรือปัญหาการเชื่อมต่อกับ OpenAI")
        return None

@st.cache_data(show_spinner=False)
def get_available_models(api_key):
    """
    ฟังก์ชันสำหรับดึงรายชื่อโมเดล GPT ที่ใช้งานได้จาก OpenAI API
    """
    if not api_key or not api_key.startswith('sk-'):
        return []
    try:
        client = openai.OpenAI(api_key=api_key)
        models = client.models.list().data
        # กรองเฉพาะโมเดลที่มี 'gpt' ในชื่อ และไม่ใช่โมเดล instruct
        # จัดเรียงโดยให้ gpt-4o มาก่อน
        gpt_models = sorted(
            [model.id for model in models if 'gpt' in model.id and 'instruct' not in model.id],
            reverse=True
        )
        if 'gpt-4o' in gpt_models:
            gpt_models.remove('gpt-4o')
            gpt_models.insert(0, 'gpt-4o')
        return gpt_models
    except Exception as e:
        # หากเกิด error (เช่น key ไม่ถูกต้อง) จะไม่แสดงข้อความ แต่จะคืนค่าเป็น list ว่าง
        print(f"ไม่สามารถดึงรายชื่อโมเดลได้: {e}")
        return []

# --- ส่วนของ UI และการทำงานของ Streamlit ---
st.set_page_config(page_title="Streamlit ChatGPT Clone", page_icon="🤖", layout="wide")

st.title("🤖 Streamlit ChatGPT Clone", help=None)
st.caption("แชทบอทอัจฉริยะที่สามารถตอบคำถามจากไฟล์ของคุณได้")

# Add upload button in a container at top right
upload_col1, upload_col2 = st.columns([1, 20])
with upload_col1:
    if st.button("📎", help="อัปโหลดไฟล์เพื่อถาม", use_container_width=True, key="top_upload_button"):
        st.session_state.show_upload_modal = True

# MongoDB connection
mongo_uri = st.secrets["MONGO_URI"]
mongo_client = MongoClient(mongo_uri)
chat_collection = mongo_client["chat_db"]["chat_logs"]

# เริ่มต้น vector_store_map
if "vector_store_map" not in st.session_state:
    st.session_state.vector_store_map = {}

# --- Sidebar สำหรับการตั้งค่า ---
with st.sidebar:
    st.header("⚙️ การตั้งค่า")
    
    # 1. รับ OpenAI API Key
    openai_api_key = st.text_input(
        "OpenAI API Key", 
        type="password",
        placeholder="กรุณาใส่ API Key ของคุณ (sk-...)",
        help="คุณสามารถหา API Key ได้จาก https://platform.openai.com/account/api-keys"
    )

    # 2. เลือกโมเดล OpenAI (แบบไดนามิก)
    available_models = get_available_models(openai_api_key)
    default_model = "gpt-4.1-nano"
    if available_models:
        if default_model in available_models:
            default_index = available_models.index(default_model)
        else:
            default_index = 0
        selected_model = st.selectbox(
            "เลือกโมเดล",
            available_models,
            index=default_index,
            help="รายชื่อโมเดลถูกดึงมาจากบัญชี OpenAI ของคุณ"
        )
    else:
        selected_model = st.selectbox(
            "เลือกโมเดล (ใส่ API Key เพื่อดูทั้งหมด)",
            ("gpt-4.1-nano", "gpt-4o-mini", "gpt-4.1-mini"),
            index=0,
            help="กรุณาใส่ API Key ที่ถูกต้องเพื่อดึงรายชื่อโมเดลทั้งหมด"
        )
    st.divider()

    st.info("แอปนี้จะทำงานใน 2 โหมด:\n1. **โหมดปกติ:** หากไม่ได้อัปโหลดและประมวลผลไฟล์ จะเป็นแชทบอททั่วไป\n2. **โหมด Deep Search:** หากประมวลผลไฟล์แล้ว บอทจะตอบคำถามโดยอ้างอิงจากเนื้อหาในไฟล์เป็นหลัก")

    # ดึงรายชื่อ chat folders ตาม user_id (คือ openai_api_key)
    chat_folder_collection = mongo_client["chat_db"]["chat_folder"]
    if openai_api_key and openai_api_key.startswith("sk-"):
        folder_doc = chat_folder_collection.find_one({"user_id": openai_api_key})
        folder_options = folder_doc.get("chat_folders", []) if folder_doc else []
        if folder_options:
            st.subheader("📂 แชทเดิมของคุณ")
            folder_names = [f["chat_name"] for f in folder_options]
            selected_folder_name = st.selectbox("เลือกแชท", folder_names, key="select_existing_chat")
            selected_folder = next((f for f in folder_options if f["chat_name"] == selected_folder_name), None)
            
            if selected_folder:
                # โหลดแชททันทีเมื่อเลือก chat folder
                if st.session_state.get("chat_id") != selected_folder["chat_folder_id"]:
                    st.session_state.messages = []
                    st.session_state.chat_id = selected_folder["chat_folder_id"]
                    st.session_state.chat_name = selected_folder["chat_name"]
                    st.session_state.system_prompt = selected_folder.get("system_prompt", "")
                    logs = chat_collection.find_one({
                        "user_id": openai_api_key,
                        "chat_folder_id": selected_folder["chat_folder_id"]
                    })
                    if logs:
                        for log in logs.get("chats", []):
                            st.session_state.messages.append({"role": "user", "content": log["question"]})
                            st.session_state.messages.append({"role": "assistant", "content": log["answer"]})
                    st.rerun()

                # --- ฟอร์มสำหรับแก้ไขชื่อแชทและ system prompt ---
                st.text_input("✏️ เปลี่ยนชื่อแชท", value=selected_folder["chat_name"], key="edit_chat_name")
                st.text_area("🧠 แก้ไข System Prompt", value=selected_folder.get("system_prompt", ""), key="edit_system_prompt")

                if st.button("💾 บันทึกการเปลี่ยนแปลง"):
                    new_name = st.session_state.get("edit_chat_name", "").strip()
                    new_prompt = st.session_state.get("edit_system_prompt", "").strip()
                    if not new_name:
                        st.warning("ชื่อแชทห้ามว่าง")
                    else:
                        chat_folder_collection.update_one(
                            {
                                "user_id": openai_api_key,
                                "chat_folders.chat_folder_id": selected_folder["chat_folder_id"]
                            },
                            {
                                "$set": {
                                    "chat_folders.$.chat_name": new_name,
                                    "chat_folders.$.system_prompt": new_prompt,
                                    "chat_folders.$.last_update": datetime.now(UTC).isoformat()
                                }
                            }
                        )
                        st.success("อัปเดตข้อมูลแชทเรียบร้อยแล้ว")
                        st.rerun()

    # --- เพิ่มส่วนเริ่มแชทใหม่ใน sidebar ---
    st.subheader("🆕 เริ่มแชทใหม่")
    if "creating_chat" not in st.session_state:
        st.session_state.creating_chat = False

    if st.button("➕ สร้างแชทใหม่"):
        st.session_state.creating_chat = True

    if st.session_state.creating_chat:
        new_name = st.text_input("📝 ชื่อแชทใหม่", key="new_chat_name_sidebar")
        new_prompt = st.text_area("🧠 system prompt (เว้นว่างได้)", key="new_system_prompt_sidebar", placeholder="กำหนด personality ของบอทที่นี่...")
        if st.button("✅ ยืนยันเริ่มแชทใหม่"):
            if not new_name.strip():
                st.warning("กรุณากรอกชื่อแชทก่อนเริ่มใช้งาน", icon="✏️")
            else:
                st.session_state.chat_id = str(uuid.uuid4())
                st.session_state.chat_name = new_name.strip()
                st.session_state.system_prompt = new_prompt.strip()
                st.session_state.messages = [{"role": "assistant", "content": "สวัสดีครับ! ผมพร้อมช่วยเหลือคุณแล้ว"}]
                st.session_state.creating_chat = False
                st.session_state.vector_store_map = {}
                st.rerun()

# --- ส่วนของการแชท ---
# --- เงื่อนไขแสดง chat input: ต้องมี chat_name ก่อนเท่านั้น ---
if "chat_name" in st.session_state and st.session_state.chat_name:
    if "chat_id" not in st.session_state:
        st.session_state.chat_id = str(uuid.uuid4())
    if "messages" not in st.session_state:
        st.session_state.messages = [{"role": "assistant", "content": "สวัสดีครับ! ผมพร้อมช่วยเหลือคุณแล้ว ใส่ API Key และอัปโหลดไฟล์ (ถ้ามี) เพื่อเริ่มใช้งานได้เลย"}]

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # ปุ่มอัปโหลดไฟล์ในหน้าหลัก (modal popup)
    if "show_upload_modal" not in st.session_state:
        st.session_state.show_upload_modal = False

    if st.session_state.show_upload_modal:
        st.markdown("### 📎 อัปโหลดไฟล์ (Deep Search Mode)")
        uploaded_file_main = st.file_uploader("เลือกไฟล์", type=["pdf", "txt", "docx", "pptx", "xlsx", "csv"], key="file_uploader_main")
        if uploaded_file_main is not None:
            if st.button("🔍 ประมวลผลไฟล์", key="process_main_file"):
                if not openai_api_key or not openai_api_key.startswith('sk-'):
                    st.warning("กรุณาใส่ OpenAI API Key ที่ถูกต้องก่อนประมวลผลไฟล์", icon="🔑")
                else:
                    with st.spinner("กำลังประมวลผลไฟล์..."):
                        raw_text = get_text_from_file(uploaded_file_main)
                        if raw_text:
                            text_chunks = get_text_chunks(raw_text)
                            vector_store = get_vector_store(text_chunks, openai_api_key)
                            if vector_store:
                                st.session_state.vector_store_map[st.session_state.chat_id] = vector_store
                                st.success("ประมวลผลไฟล์สำเร็จ! ตอนนี้คุณสามารถถามคำถามเกี่ยวกับไฟล์นี้ได้แล้ว")
                                st.session_state.show_upload_modal = False
                            else:
                                st.session_state.vector_store_map.pop(st.session_state.chat_id, None)
                        else:
                            st.error("ไม่สามารถอ่านไฟล์ได้")
                            st.session_state.vector_store_map.pop(st.session_state.chat_id, None)
        if st.button("❌ ปิด", key="close_upload_popup"):
            st.session_state.show_upload_modal = False

    # --- Chat input & attach button: moved to bottom ---
    prompt = st.chat_input("ถามคำถามของคุณ...")
    if prompt:
        if not openai_api_key or not openai_api_key.startswith('sk-'):
            st.warning("กรุณาใส่ OpenAI API Key ที่ถูกต้องในแถบด้านข้างก่อน", icon="🔑")
        else:
            openai.api_key = openai_api_key
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)

            with st.chat_message("assistant"):
                with st.spinner("กำลังคิด..."):
                    response = ""
                    chat_id = st.session_state.get("chat_id")
                    vector_store = st.session_state.vector_store_map.get(chat_id)
                    if vector_store is not None:
                        try:
                            docs = vector_store.similarity_search(query=prompt, k=3)
                            llm = OpenAI(model_name=selected_model, temperature=0.7, openai_api_key=openai_api_key)
                            memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
                            chain = load_qa_chain(llm=llm, chain_type="stuff", memory=memory)
                            response = chain.run(input_documents=docs, question=prompt)
                            st.markdown(response)
                        except Exception as e:
                            response = f"ขออภัย, เกิดข้อผิดพลาดในการค้นหาข้อมูลจากไฟล์: {e}"
                            st.error(response)
                    else:
                        try:
                            api_messages = [{"role": m["role"], "content": m["content"]} for m in st.session_state.messages]
                            # เพิ่ม system prompt
                            api_messages = [{"role": "system", "content": st.session_state.get("system_prompt", "")}] + api_messages
                            completion = openai.chat.completions.create(
                                model=selected_model,
                                messages=api_messages
                            )
                            response = completion.choices[0].message.content
                            st.markdown(response)
                        except Exception as e:
                            response = f"ขออภัย, เกิดข้อผิดพลาดในการเชื่อมต่อกับ OpenAI: {e}"
                            st.error(response)
                    st.session_state.messages.append({"role": "assistant", "content": response})
                    # Save chatlog to MongoDB (new schema)
                    user_id = openai_api_key  # ใช้ OpenAI Key เป็น user_id
                    chat_folder_id = st.session_state.get("chat_id", str(uuid.uuid4()))
                    chat_data = {
                        "question": prompt,
                        "answer": response,
                        "usage": completion.usage.to_dict() if 'completion' in locals() and hasattr(completion.usage, "to_dict") else None,
                        "timestamp": datetime.now(UTC).isoformat()
                    }
                    # upsert log เข้า chat_logs
                    chat_collection.update_one(
                        {"user_id": user_id, "chat_folder_id": chat_folder_id},
                        {
                            "$push": {"chats": chat_data},
                            "$set": {"last_update": datetime.now(UTC).isoformat()}
                        },
                        upsert=True
                    )
                    # update หรือ insert folder entry (refactored)
                    chat_folder_collection = mongo_client["chat_db"]["chat_folder"]
                    folder_doc = chat_folder_collection.find_one({"user_id": user_id})
                    folder_options = folder_doc.get("chat_folders", []) if folder_doc else []
                    folder_exists = any(f["chat_folder_id"] == chat_folder_id for f in folder_options)

                    if folder_exists:
                        # อัปเดตข้อมูลใน folder เดิม
                        chat_folder_collection.update_one(
                            {
                                "user_id": user_id,
                                "chat_folders.chat_folder_id": chat_folder_id
                            },
                            {
                                "$set": {
                                    "chat_folders.$.last_update": datetime.now(UTC).isoformat()
                                }
                            }
                        )
                    else:
                        # เพิ่ม folder ใหม่เข้าไป
                        chat_folder_collection.update_one(
                            {"user_id": user_id},
                            {
                                "$push": {
                                    "chat_folders": {
                                        "chat_folder_id": chat_folder_id,
                                        "chat_name": st.session_state.get("chat_name", "New Chat"),
                                        "system_prompt": st.session_state.get("system_prompt", ""),
                                        "model": selected_model,
                                        "create_at": datetime.now(UTC).isoformat(),
                                        "last_update": datetime.now(UTC).isoformat()
                                    }
                                }
                            },
                            upsert=True
                        )
                    st.session_state.chat_id = chat_folder_id
else:
    st.info("กรุณาเริ่มแชทใหม่โดยตั้งชื่อก่อนจึงจะสามารถเริ่มสนทนาได้")
